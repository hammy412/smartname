import argparse
import requests
import pymupdf
import fitz
from pathlib import Path
import base64
import json
import re
from docx import Document
from pptx import Presentation

def main():
    parser = argparse.ArgumentParser(description="Rename files using Ollama's LLM.")
    parser.add_argument("directory", help="Path to the directory containing files to rename.")
    parser.add_argument("--execute", action="store_true", help="Execute the renaming of files.")
    parser.add_argument("--model", default="llama3.2-vision", help="Model to use for renaming files.")
    parser.add_argument("--case", choices=["snake_case", "kebab-case", "camelCase", "PascalCase", "lowercase", "TitleCase"], default="snake_case", help="Casing style for the new filenames.")

    args = parser.parse_args()

    renamed = generate_filename_for_file(args)

    if args.execute:
        for old_path, new_name in renamed.items():
            new_name = apply_casing(new_name, args.case)
            new_path = old_path.with_name(new_name + old_path.suffix)
            old_path.rename(new_path)
            print(f"Renamed: {old_path} -> {new_path}")
    else:
        for old_path, new_name in renamed.items():
            new_name = apply_casing(new_name, args.case)
            print(f"Suggested: {old_path} -> {new_name + old_path.suffix}")


def generate_filename_for_file(args):
    directory = Path(args.directory)
    outputs = {}

    for file_path in directory.iterdir():
        if file_path.is_file():
            ext = file_path.suffix.lower()

            if ext in ['.png', '.jpg']:
                with open(file_path, "rb") as img_file:
                    image_data = base64.b64encode(img_file.read()).decode('utf-8')
                    outputs[file_path] = call_ollama_vision(image_data, args.model)
            elif ext in ['.txt', '.md', '.py']:
                with open(file_path, "r", encoding="utf-8") as text_file:
                    text = text_file.read()
                    outputs[file_path] = call_ollama_text(text, args.model)
            elif ext == '.pdf':
                doc = fitz.open(file_path)
                first_page = doc.load_page(0)
                text = first_page.get_text("text")
                if len(text.strip()) > 0:
                    outputs[file_path] = call_ollama_text(text, args.model)
                else:
                    pix = first_page.get_pixmap()
                    img_data = base64.b64encode(pix.tobytes()).decode('utf-8')
                    outputs[file_path] = call_ollama_vision(img_data, args.model)
            elif ext == '.docx':
                doc = Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs])
                outputs[file_path] = call_ollama_text(text, args.model)
            elif ext == '.pptx':
                prs = Presentation(file_path)
                text_content = []
                for slide in prs.slides:    
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content.append(shape.text)

                text = "\n".join(text_content)
                outputs[file_path] = call_ollama_text(text, args.model)

    return outputs
                
def call_ollama_vision(image_data, model):
    prompt = "Analyze this file's content and suggest a concise, 3-5 word, descriptive filename. Use underscores instead of spaces. Do not include the file extension. Respond with *only* the filename and nothing else."
    response = requests.post(
        "http://localhost:11434/api/generate",
        json={
            "model": model,
            "prompt": prompt,
            "images": [image_data],
            "options": {
                "temperature": 0.7,
                "top_p": 0.9,
                "num_predict": 50
            }
        },
        timeout=30
    )

    # Parse multiple JSON objects line by line
    suggested = ""
    for line in response.text.strip().splitlines():
        try:
            chunk = json.loads(line)
            suggested += chunk.get("response", "")
        except json.JSONDecodeError:
            continue

    suggested = sanitize_filename(suggested.strip())
    if not suggested:
        suggested = "rename_me"

    return suggested

def call_ollama_text(text, model):
    prompt = "Analyze this file's content and suggest a concise, 3-5 word, descriptive filename. Use underscores instead of spaces. Do not include the file extension. Respond with *only* the filename and nothing else."
    response = requests.post(
        "http://localhost:11434/api/generate",
        json={
            "model": model,
            "prompt": prompt + "\n\n" + text[:2000],
            "options": {
                "temperature": 0.7,
                "top_p": 0.9,
                "num_predict": 50
            }
        },
        timeout=30
    )

    # Parse multiple JSON objects line by line
    suggested = ""
    for line in response.text.strip().splitlines():
        try:
            chunk = json.loads(line)
            suggested += chunk.get("response", "")
        except json.JSONDecodeError:
            continue

    suggested = sanitize_filename(suggested.strip())
    if not suggested:
        suggested = "rename_me"

    return suggested

def apply_casing(filename, style="snake_case"):
     # Normalize filename into a list of lowercase words
    words = re.split(r'[_\-\s]+', filename.strip())
    words = [w for w in words if w]  # remove empties
    if not words:
        return "rename_me"

    style = style.lower()

    if style == "snake_case":
        return "_".join(w.lower() for w in words)

    elif style == "kebab-case":
        return "-".join(w.lower() for w in words)

    elif style == "camelcase":
        return words[0].lower() + "".join(w.capitalize() for w in words[1:])

    elif style == "pascalcase":
        return "".join(w.capitalize() for w in words)

    elif style == "lowercase":
        return " ".join(w.lower() for w in words)

    elif style == "titlecase":
        return " ".join(w.capitalize() for w in words)

    else:
        return "_".join(w.lower() for w in words)  # default fallback

def sanitize_filename(filename):

    # Remove illegal characters for filenames on most OS
    filename = re.sub(r'[\\/*?:"<>|]', '', filename)
    
    # Remove surrounding whitespace
    filename = filename.strip()
    
    # Remove any existing extension like .txt, .jpg, etc.
    filename = re.sub(r'\.[a-zA-Z0-9]{1,5}$', '', filename)
    
    # Remove non-word punctuation (except spaces, underscores, and hyphens)
    filename = re.sub(r'[^\w\s\-_]', '', filename)
    
    # Collapse multiple spaces, underscores, or hyphens into single spaces (for better casing)
    filename = re.sub(r'[\s_\-]+', ' ', filename)
    
    # Remove extra leading/trailing spaces again
    filename = filename.strip()

    # Ensure not empty
    if not filename:
        filename = "rename_me"

    return filename




if __name__ == "__main__":
    main()